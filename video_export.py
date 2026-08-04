# -*- coding: utf-8 -*-
"""
video_export.py — 短视频批改导出（每生 PDF / 全班 Excel / 讲评 PPT）
════════════════════════════════════════════════════════════════
- 字体：fonts/NotoSansSC-Regular.ttf + Bold（与教师版同一套，随仓库分发）
- 分数展示决策（2026-07-20）：短视频评分表是刘老师自订规则，非 SEAB
  校准资产，因此 PDF/Excel 直接展示各维度分与总分（与作文"只给等级"不同）。
- PPT 定位是"讲评底稿"，可编辑 pptx；素材来自 REVIEW_AGG 聚合结果
  （人名已在提示词层匿名化）。
- 生成失败绝不影响页面 — 调用处必须 try/except。
"""

import io
import json
import os
import zipfile
from datetime import datetime, timezone, timedelta

from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.lib.colors import HexColor
from reportlab.lib.styles import ParagraphStyle
from reportlab.platypus import (BaseDocTemplate, Frame, PageTemplate,
                                Paragraph, Spacer, Table, TableStyle)
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

from openpyxl import Workbook
from openpyxl.styles import Font as XFont, PatternFill, Alignment

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from video_rubric import (dims, total_max, level_of,
                          total_level_of, member_mode)
from video_engine import total_of, review_flags, MEMBER_MARKS

MEMBER_MARK_LABELS = MEMBER_MARKS

EXPORT_VERSION = "video-2.0-20260804"

NAVY = HexColor("#1F4E79")
GOLD = HexColor("#f0c27f")
TEXT = HexColor("#2c2c2a")
MUTED = HexColor("#888888")
GRAY_BG = HexColor("#f5f4f0")
RED = HexColor("#c62828")
GREEN = HexColor("#2e7d32")

_FONTS_REGISTERED = False


def _register_fonts():
    global _FONTS_REGISTERED
    if _FONTS_REGISTERED:
        return
    base = os.path.join(os.path.dirname(os.path.abspath(__file__)), "fonts")
    reg = os.path.join(base, "NotoSansSC-Regular.ttf")
    bold = os.path.join(base, "NotoSansSC-Bold.ttf")
    if not os.path.exists(reg):
        raise RuntimeError(f"缺少中文字体文件: {reg}")
    pdfmetrics.registerFont(TTFont("NotoSC", reg))
    pdfmetrics.registerFont(TTFont("NotoSC-Bold",
                                   bold if os.path.exists(bold) else reg))
    _FONTS_REGISTERED = True


def _sg_today():
    return datetime.now(timezone(timedelta(hours=8))).strftime("%Y-%m-%d")


def _styles():
    def s(name, **kw):
        base = dict(fontName="NotoSC", fontSize=9.5, leading=16, textColor=TEXT)
        base.update(kw)
        return ParagraphStyle(name, **base)
    return {
        "title": s("title", fontName="NotoSC-Bold", fontSize=15, leading=22,
                   textColor=NAVY),
        "sub": s("sub", fontSize=8.5, leading=13, textColor=MUTED),
        "sec": s("sec", fontName="NotoSC-Bold", fontSize=11, leading=17,
                 textColor=NAVY, spaceBefore=8),
        "body": s("body"),
        "big": s("big", fontName="NotoSC-Bold", fontSize=13, leading=20,
                 textColor=NAVY),
        "small": s("small", fontSize=8.6, leading=14),
        "quote": s("quote", fontSize=8.6, leading=14,
                   textColor=HexColor("#555555"), leftIndent=6),
    }


def _esc(t):
    return (str(t or "").replace("&", "&amp;").replace("<", "&lt;")
            .replace(">", "&gt;"))


# ─────────────────────────────────────────────────────────────
# 每生 PDF
# ─────────────────────────────────────────────────────────────

def build_student_pdf(item, class_name, topic, rubric,
                      sid=None, member=None):
    """学生报告 PDF。

    2026-08-04 改版（刘老师要求）：
    - 抬头改成「第 NN 号视频」，全文不出现"学号"字样——一个视频里有好几位
      同学，按学号称呼会让学生以为整份报告在说他一个人。
    - B 案版式：总分与等级 + 四维度得分小表 + 2条优点 + 2条不足。
      去掉 evidence 逐字引用、issues 三段式、口语小数据——那些是给老师
      复核用的，学生看不完也用不上。
    - 【隐私分层】组共享层（总分/等级/维度表/优点/不足）只评作品，绝不点名
      个人；点名的内容全部放进"你个人的部分"，且只出现在该生自己那一份里。
      理由：四五个人看同一份报告，写"某位同学在念稿"等于当众点名。
    sid/member 为空时（个人作业模板）自动退化成原来的单人报告。
    """
    _register_fonts()
    st = _styles()
    ai = item.get("ai_result") or {}
    dims_r = ai.get("dimensions", {})
    fs = item.get("final_scores") or {}
    total = total_of(fs)
    DIMS = dims(rubric)
    TM = total_max(rubric)
    tlv = total_level_of(rubric, total)

    buf = io.BytesIO()
    doc = BaseDocTemplate(buf, pagesize=A4,
                          leftMargin=18 * mm, rightMargin=18 * mm,
                          topMargin=16 * mm, bottomMargin=16 * mm)
    frame = Frame(doc.leftMargin, doc.bottomMargin, doc.width, doc.height)
    doc.addPageTemplates([PageTemplate(frames=[frame])])
    story = []

    # ── 抬头：按视频编号称呼，不出现学号 ──
    vid_no = item.get("item_key") or ""
    head = f"第 {vid_no} 号视频　小组作业批改报告" if vid_no else "短视频作业批改报告"
    story.append(Paragraph(_esc(head), st["title"]))
    story.append(Paragraph(
        f"{_esc(class_name)} ｜ 题目 {_esc(topic)} ｜ {_sg_today()}", st["sub"]))
    story.append(Spacer(1, 5 * mm))

    # ── 成绩：总分 + 等级 ──
    if tlv:
        story.append(Paragraph(
            f"<b>小组成绩：{total} / {TM}　等级 {_esc(tlv['grade'])}"
            f"（{_esc(tlv['label'])}）</b>", st["big"]))
    else:
        story.append(Paragraph(f"<b>总分：{total} / {TM}</b>", st["big"]))
    story.append(Spacer(1, 3 * mm))

    # ── 四维度得分小表（B 案）──
    rows = [["评分项", "得分", "等级"]]
    for d in DIMS:
        sc = fs.get(d["key"], 0)
        lv = level_of(d, sc)
        rows.append([d["name"], f"{sc} / {d['max']}",
                     f"{lv['grade']} {lv['label']}"])
    tbl = Table(rows, colWidths=[62 * mm, 32 * mm, 40 * mm])
    tbl.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (-1, -1), "NotoSC"),
        ("FONTNAME", (0, 0), (-1, 0), "NotoSC-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 9.5),
        ("BACKGROUND", (0, 0), (-1, 0), GRAY_BG),
        ("GRID", (0, 0), (-1, -1), 0.4, HexColor("#cccccc")),
        ("ALIGN", (1, 0), (-1, -1), "CENTER"),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    story.append(tbl)
    story.append(Spacer(1, 4 * mm))

    # ── 优点（组共享层，只评作品）──
    strengths = [x for x in (ai.get("report_strengths") or []) if str(x).strip()]
    if not strengths:                      # 兼容旧批改记录
        strengths = [x for x in [ai.get("top_strength")] if x]
    if strengths:
        story.append(Paragraph("做得好的地方", st["sec"]))
        for x in strengths[:3]:
            story.append(Paragraph(f"• {_esc(x)}", st["body"]))

    # ── 不足 + 怎么改 ──
    imps = [x for x in (ai.get("report_improvements") or [])
            if isinstance(x, dict) and str(x.get("what", "")).strip()]
    if not imps and ai.get("top_issue"):   # 兼容旧批改记录
        imps = [{"what": ai["top_issue"],
                 "how": ai.get("next_level_advice", "")}]
    if imps:
        story.append(Paragraph("可以再改进的地方", st["sec"]))
        for x in imps[:3]:
            line = _esc(x.get("what", ""))
            if str(x.get("how", "")).strip():
                line += f"　→　{_esc(x['how'])}"
            story.append(Paragraph(f"• {line}", st["body"]))

    # ── 个人私有层（只出现在该生自己那一份）──
    if member:
        story.append(Spacer(1, 3 * mm))
        story.append(Paragraph("你个人的部分", st["sec"]))
        story.append(Paragraph(
            f"<b>你的得分：{member.get('score', '')} / {TM}</b>", st["body"]))
        if member.get("note"):
            story.append(Paragraph(_esc(member["note"]), st["small"]))
        if member.get("teacher_note"):
            story.append(Paragraph(_esc(member["teacher_note"]), st["body"]))

    # ── 老师的话 ──
    if item.get("teacher_comment"):
        story.append(Paragraph("老师的话", st["sec"]))
        story.append(Paragraph(_esc(item["teacher_comment"]), st["body"]))

    story.append(Spacer(1, 6 * mm))
    story.append(Paragraph("AI 批改 + 老师复核", st["sub"]))
    doc.build(story)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────
# 全班 Excel
# ─────────────────────────────────────────────────────────────

def build_class_excel(items, class_name, topic, rubric):
    """一行一个学生（两人组各占一行）。返回 bytes。"""
    wb = Workbook()
    ws = wb.active
    ws.title = "短视频成绩总表"
    DIMS = dims(rubric)
    TM = total_max(rubric)
    MM = member_mode(rubric)
    # 2026-08-04：小组模板多两列——个人分才是最终要录进成绩册的分数
    header = (["学号", "视频编号"]
              + [f"{d['name']}/{d['max']}" for d in DIMS]
              + [f"组分/{TM}"]
              + ([f"个人分/{TM}", "档内位置"] if MM else [])
              + ["主要问题", "突出优点",
                 "距上一档建议", "复核提示", "教师最终评分（手写录入）"])
    ws.append([f"{class_name} 短视频作业 ｜ 题目：{topic} ｜ {_sg_today()}"])
    ws.append(header)
    yellow = PatternFill("solid", start_color="FFF3CD")
    for c in ws[2]:
        c.font = XFont(bold=True)
    for it in items:
        ai = it.get("ai_result") or {}
        fs = it.get("final_scores") or {}
        flags = review_flags(it, rubric)
        ids = it.get("student_ids") or []
        key = it.get("item_key") or ""
        ms = it.get("member_scores") or {}
        if isinstance(ms, str):
            ms = json.loads(ms)
        for sid in ids:
            mem = ms.get(sid) or {}
            row = ([sid, key]
                   + [fs.get(d["key"], "") for d in DIMS]
                   + [total_of(fs)]
                   + ([mem.get("score", ""),
                       MEMBER_MARK_LABELS.get(mem.get("mark"), "")]
                      if MM else [])
                   + [ai.get("top_issue", ""), ai.get("top_strength", ""),
                      ai.get("next_level_advice", ""),
                      "；".join(flags), ""])
            ws.append(row)
            if flags:
                for c in ws[ws.max_row]:
                    c.fill = yellow
    widths = ([8, 10] + [12] * len(DIMS) + [9]
              + ([9, 16] if MM else []) + [22, 22, 30, 26, 20])
    from openpyxl.utils import get_column_letter
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w
    for row in ws.iter_rows(min_row=3):
        for c in row:
            c.alignment = Alignment(vertical="top", wrap_text=True)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────
# 讲评 PPT（讲评底稿，可编辑）
# ─────────────────────────────────────────────────────────────

_PPT_NAVY = RGBColor(0x1F, 0x4E, 0x79)
_PPT_RED = RGBColor(0xC6, 0x28, 0x28)
_PPT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
_PPT_DARK = RGBColor(0x2C, 0x2C, 0x2A)


def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式


def _txt(slide, x, y, w, h, text, size, bold=False, color=_PPT_DARK,
         align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in str(text).split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        if align:
            p.alignment = align
    return box


def build_review_ppt(agg, class_name, topic, grade_dist):
    """agg: REVIEW_AGG 聚合 JSON；grade_dist: {'A':n,...} 按总分折算的分布。
    字号按课堂投影标准放大（沿用作文 PPT 经验）。返回 bytes。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 封面
    s = _slide(prs)
    _txt(s, 1.2, 2.2, 11, 1.2, f"《{topic}》短视频作业讲评", 44, True, _PPT_NAVY)
    _txt(s, 1.2, 3.6, 11, 0.8, f"{class_name} ｜ {_sg_today()} ｜ AI 辅助整理 · 老师复核",
         20, False, _PPT_DARK)

    # 总览
    s = _slide(prs)
    _txt(s, 0.8, 0.5, 11, 0.9, "全班总体表现", 36, True, _PPT_NAVY)
    dist_line = "　".join(f"{g}档 {n} 份" for g, n in grade_dist.items() if n)
    _txt(s, 0.8, 1.7, 11.5, 1.0, dist_line, 24, True, _PPT_DARK)
    _txt(s, 0.8, 2.8, 11.5, 3.5, agg.get("overview", ""), 24)

    # 共性问题：问题 → 例子 → 为什么错 → 怎么改（一题一页）
    for i, iss in enumerate(agg.get("common_issues", [])[:4], 1):
        s = _slide(prs)
        _txt(s, 0.8, 0.4, 11.5, 0.9,
             f"共性问题 {i}：{iss.get('title', '')}（{iss.get('count_hint', '')}）",
             32, True, _PPT_RED)
        y = 1.5
        for ex in iss.get("examples", [])[:2]:
            _txt(s, 1.0, y, 11.2, 0.85,
                 f"“{ex.get('quote', '')}”  [{ex.get('t', '')}]", 22,
                 color=RGBColor(0x55, 0x55, 0x55))
            y += 0.95
        _txt(s, 0.8, y + 0.1, 11.5, 1.4,
             f"为什么错：{iss.get('why', '')}", 24)
        _txt(s, 0.8, y + 1.7, 11.5, 1.8,
             f"怎么改：{iss.get('how', '')}", 24, True, _PPT_GREEN)

    # 佳句欣赏
    s = _slide(prs)
    _txt(s, 0.8, 0.4, 11.5, 0.9, "本班佳句欣赏", 36, True, _PPT_GREEN)
    y = 1.5
    for p in agg.get("praise", [])[:6]:
        _txt(s, 1.0, y, 11.3, 0.9,
             f"“{p.get('quote', '')}” —— {p.get('why', '')}", 22)
        y += 0.95

    # 下次拍摄前检查清单
    s = _slide(prs)
    _txt(s, 0.8, 0.4, 11.5, 0.9, "下次拍摄前，检查这五件事", 36, True, _PPT_NAVY)
    y = 1.6
    for i, c in enumerate(agg.get("checklist", [])[:5], 1):
        _txt(s, 1.0, y, 11.3, 0.9, f"☑ {c}", 26)
        y += 0.95

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────
# 打包
# ─────────────────────────────────────────────────────────────

def grade_dist_of(items, rubric):
    """按总分占满分比例折四档（A≥80% B≥60% C≥40% D<40%），仅讲评总览用。"""
    tm = total_max(rubric) or 1
    dist = {"A": 0, "B": 0, "C": 0, "D": 0}
    for it in items:
        p = total_of(it.get("final_scores")) / tm
        g = "A" if p >= 0.8 else "B" if p >= 0.6 else "C" if p >= 0.4 else "D"
        dist[g] += 1
    return dist


def build_student_zip(items, class_name, topic, rubric):
    """【发给学生的包】只有 PDF，一人一份。

    2026-08-04 决策：与老师自留包彻底分开。原来 Excel（全班分数）和讲评 PPT
    跟学生 PDF 混在同一个 ZIP 里，老师手一滑整包传上谷歌课室，全班分数就
    公开了。分成两个包是防呆，不是洁癖。
    小组模板：每位组员一份，文件名「组号_学号.pdf」，前半段（作品评价）
    四五个人完全相同，后半段「你个人的部分」各不相同。
    """
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for it in items:
            ids = it.get("student_ids") or []
            key = it.get("item_key") or (ids[0] if ids else "unknown")
            ms = it.get("member_scores") or {}
            if isinstance(ms, str):
                ms = json.loads(ms)
            if member_mode(rubric) and ms:
                for sid, mem in ms.items():
                    pdf = build_student_pdf(it, class_name, topic, rubric,
                                            sid=sid, member=mem)
                    zf.writestr(f"{key}_{sid}.pdf", pdf)
            else:
                pdf = build_student_pdf(it, class_name, topic, rubric)
                for sid in (ids or [key]):
                    zf.writestr(f"{sid}.pdf", pdf)
    return buf.getvalue()


def build_teacher_zip(items, class_name, topic, rubric, agg=None):
    """【老师自留的包】全班 Excel + 讲评 PPT。★绝不可发给学生★"""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(f"{class_name}_短视频成绩总表.xlsx",
                    build_class_excel(items, class_name, topic, rubric))
        if agg:
            zf.writestr(f"{class_name}_短视频讲评.pptx",
                        build_review_ppt(agg, class_name, topic,
                                         grade_dist_of(items, rubric)))
    return buf.getvalue()
